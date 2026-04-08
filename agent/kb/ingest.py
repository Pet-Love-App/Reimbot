from __future__ import annotations

import argparse
import json
import re
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, Iterable, List, Tuple
from docx import Document
from pptx import Presentation


SUPPORTED_SUFFIXES = {".txt", ".md", ".docx", ".pptx", ".xlsx", ".xls", ".pdf"}
_SENTENCE_BOUNDARY_RE = re.compile(r"(?<=[。！？!?；;])")


def _clean_text(text: str) -> str:
    cleaned = text.replace("\u3000", " ")
    cleaned = re.sub(r"\r\n?", "\n", cleaned)
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)
    cleaned = re.sub(r"[ \t]{2,}", " ", cleaned)
    return cleaned.strip()


def _read_text_file(path: Path) -> str:
    try:
        return path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        return path.read_text(encoding="gbk", errors="ignore")


def _read_docx(path: Path) -> str:
    doc = Document(path)
    parts: List[str] = []

    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if text:
            parts.append(text)

    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))

    return "\n".join(parts)


def _read_pptx(path: Path) -> str:
    prs = Presentation(path)
    parts: List[str] = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_texts: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = str(shape.text).strip()
                if text:
                    slide_texts.append(text)
        if slide_texts:
            parts.append(f"[Slide {slide_idx}]\n" + "\n".join(slide_texts))

    return "\n\n".join(parts)


def _read_excel(path: Path) -> str:
    try:
        import pandas as pd
    except Exception:
        return ""
    xls = pd.ExcelFile(path)
    blocks: List[str] = []

    for sheet_name in xls.sheet_names:
        df = xls.parse(sheet_name=sheet_name, dtype=str).fillna("")
        if df.empty:
            continue
        sheet_text = df.astype(str).apply(lambda row: " | ".join(row.tolist()), axis=1).tolist()
        blocks.append(f"[Sheet {sheet_name}]\n" + "\n".join(sheet_text))

    return "\n\n".join(blocks)


def _read_pdf(path: Path) -> str:
    try:
        import fitz  # PyMuPDF
    except Exception:
        return ""

    parts: List[str] = []
    try:
        with fitz.open(path) as doc:
            for page in doc:
                text = (page.get_text("text") or "").strip()
                if text:
                    parts.append(text)
    except Exception:
        return ""
    return "\n\n".join(parts).strip()


def _extract_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix in {".txt", ".md"}:
        return _read_text_file(path)
    if suffix == ".docx":
        return _read_docx(path)
    if suffix == ".pptx":
        return _read_pptx(path)
    if suffix in {".xlsx", ".xls"}:
        return _read_excel(path)
    if suffix == ".pdf":
        return _read_pdf(path)
    return ""


def _split_chunks(text: str, *, chunk_size: int, overlap: int) -> List[str]:
    normalized = _clean_text(text)
    if not normalized:
        return []

    paragraphs = [para.strip() for para in normalized.split("\n\n") if para.strip()]
    units: List[str] = []

    def _split_by_sentence(para: str) -> List[str]:
        if len(para) <= chunk_size:
            return [para]
        raw_parts = [part.strip() for part in _SENTENCE_BOUNDARY_RE.split(para) if part.strip()]
        if not raw_parts:
            raw_parts = [para]

        refined: List[str] = []
        for part in raw_parts:
            if len(part) <= chunk_size:
                refined.append(part)
                continue
            start = 0
            while start < len(part):
                end = min(start + chunk_size, len(part))
                piece = part[start:end].strip()
                if piece:
                    refined.append(piece)
                if end >= len(part):
                    break
                start = end
        return refined

    for para in paragraphs:
        units.extend(_split_by_sentence(para))

    chunks: List[str] = []
    current_units: List[str] = []
    current_len = 0

    for unit in units:
        separator_len = 2 if current_units else 0
        projected_len = current_len + separator_len + len(unit)
        if projected_len <= chunk_size:
            if current_units:
                current_len += 2
            current_units.append(unit)
            current_len += len(unit)
            continue

        if current_units:
            chunks.append("\n\n".join(current_units).strip())

        if overlap > 0 and chunks:
            tail = chunks[-1][-overlap:].strip()
            current_units = [tail, unit] if tail else [unit]
            current_len = len("\n\n".join(current_units))
        else:
            current_units = [unit]
            current_len = len(unit)

    if current_units:
        chunks.append("\n\n".join(current_units).strip())

    return [chunk for chunk in chunks if chunk]


def _iter_files(source_dir: Path) -> Iterable[Path]:
    for path in source_dir.rglob("*"):
        if path.is_file() and path.suffix.lower() in SUPPORTED_SUFFIXES:
            yield path


def _infer_category(relative_source: str) -> Tuple[str, str]:
    parts = [part for part in relative_source.replace("\\", "/").split("/") if part]
    if not parts:
        return "未分类", ""
    if len(parts) == 1:
        return parts[0], ""
    return parts[0], "/".join(parts[1:-1])


def _safe_relative(path_str: str, root: Path) -> str:
    path = Path(path_str)
    try:
        return path.resolve().relative_to(root.resolve()).as_posix()
    except Exception:
        return path.name


def _build_chunk_id(relative_source: str, index: int) -> str:
    return f"{relative_source}#{index}"


def _extract_parsed_text(parsed_dir: Path, fallback_doc: Dict[str, Any]) -> str:
    text_md = parsed_dir / "text.md"
    if text_md.exists():
        text = text_md.read_text(encoding="utf-8", errors="ignore").strip()
        if text:
            return text

    parts: List[str] = []
    for section in fallback_doc.get("sections", []) or []:
        heading = str(section.get("heading", "")).strip()
        body = str(section.get("text", "")).strip()
        block = "\n".join([item for item in [heading, body] if item]).strip()
        if block:
            parts.append(block)
    for slide in fallback_doc.get("slides", []) or []:
        title = str(slide.get("title", "")).strip()
        body = str(slide.get("text", "")).strip()
        block = "\n".join([item for item in [title, body] if item]).strip()
        if block:
            parts.append(block)
    for table in fallback_doc.get("tables", []) or []:
        headers = table.get("headers", []) or []
        if headers:
            parts.append(" | ".join(str(h) for h in headers if str(h).strip()))
    return "\n\n".join(parts).strip()


def _build_chunks_from_text(
    *,
    text: str,
    file_path: Path,
    source_dir: Path,
    title: str,
    chunk_size: int,
    overlap: int,
) -> List[Dict[str, Any]]:
    relative_source = _safe_relative(str(file_path), source_dir)
    category, subcategory = _infer_category(relative_source)
    file_chunks = _split_chunks(text, chunk_size=chunk_size, overlap=overlap)
    if not file_chunks:
        return []

    doc_type = file_path.suffix.lstrip(".").lower()
    built: List[Dict[str, Any]] = []
    for idx, chunk in enumerate(file_chunks, start=1):
        built.append(
            {
                "id": _build_chunk_id(relative_source, idx),
                "source": relative_source,
                "title": f"{title}-片段{idx}",
                "content": chunk,
                "category": category,
                "subcategory": subcategory,
                "doc_type": doc_type,
            }
        )
    return built


def _collect_chunks_from_parse_results(
    parse_result: Dict[str, Any],
    *,
    source_dir: Path,
    chunk_size: int,
    overlap: int,
) -> Tuple[int, List[Dict[str, Any]]]:
    chunks: List[Dict[str, Any]] = []
    file_count = 0

    for item in parse_result.get("results", []):
        if not isinstance(item, dict):
            continue
        if item.get("status") == "error":
            continue
        file_path = str(item.get("file_path", "")).strip()
        parsed_dir_str = str(item.get("parsed_dir", "")).strip()
        if not file_path or not parsed_dir_str:
            continue

        parsed_dir = Path(parsed_dir_str)
        doc_json_path = parsed_dir / "document.json"
        if not doc_json_path.exists():
            continue

        document_json = json.loads(doc_json_path.read_text(encoding="utf-8", errors="ignore"))
        source = document_json.get("source", {}) or {}
        file_type = str(source.get("file_type", Path(file_path).suffix.lstrip(".").lower()))
        title = str(document_json.get("title", "")).strip() or Path(file_path).stem
        relative_source = _safe_relative(file_path, source_dir)
        category, subcategory = _infer_category(relative_source)

        text = _extract_parsed_text(parsed_dir, document_json)
        file_chunks = _split_chunks(text, chunk_size=chunk_size, overlap=overlap)
        if not file_chunks:
            continue

        file_count += 1
        for idx, chunk in enumerate(file_chunks, start=1):
            chunks.append(
                {
                    "id": _build_chunk_id(relative_source, idx),
                    "source": relative_source,
                    "title": f"{title}-片段{idx}",
                    "content": chunk,
                    "category": category,
                    "subcategory": subcategory,
                    "doc_type": file_type,
                }
            )

    return file_count, chunks


def _persist_to_chroma(chunks: List[Dict[str, Any]], output_file: Path) -> None:
    try:
        import chromadb
        from chromadb.utils import embedding_functions

        db_path = output_file.parent / "chroma_db"
        client = chromadb.PersistentClient(path=str(db_path))

        class JinaEmbeddingFunction(embedding_functions.EmbeddingFunction):
            def __call__(self, input: chromadb.Documents) -> chromadb.Embeddings:
                from agent.kb.retriever import _get_model

                model = _get_model()
                if model is None:
                    raise RuntimeError("Embedding model is not available")
                embeddings = model.encode(input, convert_to_numpy=True, show_progress_bar=False, batch_size=32)
                return embeddings.tolist()

        emb_fn = JinaEmbeddingFunction()
        collection = client.get_or_create_collection(
            name="reimbursement_kb",
            embedding_function=emb_fn,
        )

        if chunks:
            ids = [str(c["id"]) for c in chunks]
            documents = [(str(c.get("title", "")) + "\n" + str(c.get("content", ""))).strip() for c in chunks]
            metadatas = [
                {
                    "source": str(c.get("source", "")),
                    "title": str(c.get("title", "")),
                    "content": str(c.get("content", "")),
                    "category": str(c.get("category", "")),
                    "subcategory": str(c.get("subcategory", "")),
                    "doc_type": str(c.get("doc_type", "")),
                }
                for c in chunks
            ]
            collection.delete(ids=ids)
            collection.upsert(documents=documents, metadatas=metadatas, ids=ids)
    except ImportError:
        pass


def _write_payload(output_file: Path, source_dir: Path, chunk_size: int, overlap: int, file_count: int, chunks: List[Dict[str, Any]], *, strategy: str) -> None:
    category_counts: Dict[str, int] = {}
    for chunk in chunks:
        category = str(chunk.get("category", "未分类")) or "未分类"
        category_counts[category] = category_counts.get(category, 0) + 1

    payload = {
        "metadata": {
            "source_dir": str(source_dir),
            "built_at": datetime.now().isoformat(timespec="seconds"),
            "chunk_size": chunk_size,
            "overlap": overlap,
            "file_count": file_count,
            "chunk_count": len(chunks),
            "strategy": strategy,
            "category_counts": category_counts,
        },
        "chunks": chunks,
    }

    output_file.parent.mkdir(parents=True, exist_ok=True)
    output_file.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")


def build_kb(
    source_dir: Path,
    output_file: Path,
    *,
    chunk_size: int,
    overlap: int,
    persist_chroma: bool,
) -> Tuple[int, int]:
    chunks: List[Dict[str, str]] = []
    file_count = 0

    for file_path in _iter_files(source_dir):
        text = _extract_text(file_path)
        file_chunks = _split_chunks(text, chunk_size=chunk_size, overlap=overlap)
        if not file_chunks:
            continue

        file_count += 1
        relative_source = file_path.relative_to(source_dir).as_posix()

        for idx, chunk in enumerate(file_chunks, start=1):
            chunks.append(
                {
                    "id": f"{relative_source}#{idx}",
                    "source": relative_source,
                    "title": f"{file_path.stem}-片段{idx}",
                    "content": chunk,
                }
            )

    _write_payload(
        output_file=output_file,
        source_dir=source_dir,
        chunk_size=chunk_size,
        overlap=overlap,
        file_count=file_count,
        chunks=chunks,
        strategy="direct",
    )
    if persist_chroma:
        _persist_to_chroma(chunks, output_file)

    return file_count, len(chunks)


def build_kb_with_parser(
    source_dir: Path,
    output_file: Path,
    *,
    chunk_size: int,
    overlap: int,
    parse_output_dir: Path,
    persist_chroma: bool,
    kb_name: str = "finance_documents",
    parser_suffixes: Tuple[str, ...] = (".docx", ".pptx", ".xlsx", ".xls", ".pdf", ".md", ".markdown", ".txt"),
) -> Tuple[int, int]:
    from agent.parser.main import parse_single_file

    suffix_set = {s.lower() for s in parser_suffixes}
    parse_output_dir.mkdir(parents=True, exist_ok=True)
    parse_result: Dict[str, Any] = {"results": []}
    fallback_files: List[Path] = []

    for path in sorted(source_dir.rglob("*"), key=lambda p: p.name):
        if not path.is_file():
            continue
        suffix = path.suffix.lower()
        if suffix not in suffix_set and suffix not in SUPPORTED_SUFFIXES:
            continue
        if suffix in suffix_set:
            try:
                item_result = parse_single_file(
                    file_path=str(path),
                    parsed_output_dir=str(parse_output_dir),
                    kb_name=kb_name,
                )
                parse_result["results"].append(item_result)
                if str(item_result.get("status", "")).lower() == "error":
                    fallback_files.append(path)
            except Exception:
                fallback_files.append(path)
            continue
        fallback_files.append(path)

    file_count, chunks = _collect_chunks_from_parse_results(
        parse_result,
        source_dir=source_dir,
        chunk_size=chunk_size,
        overlap=overlap,
    )

    chunked_sources = {str(item.get("source", "")) for item in chunks}
    fallback_added = 0
    for path in fallback_files:
        relative_source = _safe_relative(str(path), source_dir)
        if relative_source in chunked_sources:
            continue
        text = _extract_text(path)
        direct_chunks = _build_chunks_from_text(
            text=text,
            file_path=path,
            source_dir=source_dir,
            title=path.stem,
            chunk_size=chunk_size,
            overlap=overlap,
        )
        if not direct_chunks:
            continue
        chunks.extend(direct_chunks)
        chunked_sources.add(relative_source)
        fallback_added += 1

    file_count += fallback_added

    _write_payload(
        output_file=output_file,
        source_dir=source_dir,
        chunk_size=chunk_size,
        overlap=overlap,
        file_count=file_count,
        chunks=chunks,
        strategy="parser",
    )
    if persist_chroma:
        _persist_to_chroma(chunks, output_file)
    return file_count, len(chunks)


def main() -> None:
    parser = argparse.ArgumentParser(description="构建报销知识库索引")
    parser.add_argument("--source", default="docs/reimbursement", help="资料目录")
    parser.add_argument("--output", default="data/kb/reimbursement_kb.json", help="索引输出路径")
    parser.add_argument("--chunk-size", type=int, default=700, help="单片段最大字符数")
    parser.add_argument("--overlap", type=int, default=100, help="分片重叠字符数")
    parser.add_argument(
        "--strategy",
        choices=["direct", "parser"],
        default="direct",
        help="direct: 直接读取文档内容；parser: 先用 parser 解析后再入库（推荐 docs/documents）",
    )
    parser.add_argument(
        "--parsed-output",
        default="docs/parsed/documents",
        help="parser 策略下的解析产物目录",
    )
    parser.add_argument("--kb-name", default="finance_documents", help="parser 策略对应的知识库名")
    parser.add_argument(
        "--parser-suffixes",
        default=".docx,.pptx,.xlsx,.xls,.pdf,.md,.markdown,.txt",
        help="parser 策略参与解析的后缀（逗号分隔，未包含的支持格式会自动尝试 direct 降级）",
    )
    parser.add_argument(
        "--persist-chroma",
        action="store_true",
        help="同时写入 Chroma 向量库（当前环境若存在 numpy/chromadb 崩溃可不启用）",
    )

    args = parser.parse_args()

    source_dir = Path(args.source).resolve()
    output_file = Path(args.output).resolve()

    if not source_dir.exists() or not source_dir.is_dir():
        raise FileNotFoundError(f"资料目录不存在: {source_dir}")

    chunk_size = max(args.chunk_size, 200)
    overlap = max(min(args.overlap, args.chunk_size - 1), 0)
    if args.strategy == "parser":
        parser_suffixes = tuple(
            s.strip().lower() if s.strip().startswith(".") else "." + s.strip().lower()
            for s in args.parser_suffixes.split(",")
            if s.strip()
        )
        file_count, chunk_count = build_kb_with_parser(
            source_dir=source_dir,
            output_file=output_file,
            chunk_size=chunk_size,
            overlap=overlap,
            parse_output_dir=Path(args.parsed_output).resolve(),
            persist_chroma=args.persist_chroma,
            kb_name=args.kb_name,
            parser_suffixes=parser_suffixes,
        )
    else:
        file_count, chunk_count = build_kb(
            source_dir=source_dir,
            output_file=output_file,
            chunk_size=chunk_size,
            overlap=overlap,
            persist_chroma=args.persist_chroma,
        )

    print(
        json.dumps(
            {
                "ok": True,
                "source": str(source_dir),
                "output": str(output_file),
                "file_count": file_count,
                "chunk_count": chunk_count,
                "strategy": args.strategy,
            },
            ensure_ascii=False,
            indent=2,
        )
    )


if __name__ == "__main__":
    main()
